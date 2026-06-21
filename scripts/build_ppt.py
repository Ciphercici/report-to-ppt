"""
PPT generation script for the report-to-ppt skill.

Three modes:
  Outline: python build_ppt.py outline.slide.md [-o output.pptx]
           Parses Marp-format .slide.md (hand-crafted by Claude) and
           generates well-designed PPTX. This is the quality path:
           Claude does content design, script does formatting.

  Auto:    python build_ppt.py report.md [-o output.pptx]
           Parses raw report markdown and auto-generates PPTX.
           Quick draft — less polished than outline mode.

  Manual:  Edit slide functions below, then python build_ppt.py
           Full control over every slide. For fine-tuning.
"""
import os, re, sys
from pathlib import Path
from urllib.parse import unquote
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ══════════════════════════════════════════
#  DESIGN SYSTEM
# ══════════════════════════════════════════

W = Inches(13.333)
H = Inches(7.5)

CLR = {
    "bg":       RGBColor(0xFF, 0xFF, 0xFF),
    "dark":     RGBColor(0x1A, 0x1A, 0x1A),
    "body":     RGBColor(0x44, 0x44, 0x44),
    "muted":    RGBColor(0x88, 0x88, 0x88),
    "accent":   RGBColor(0x3B, 0x82, 0xC4),
    "accent2":  RGBColor(0xE8, 0x4D, 0x4D),
    "light":    RGBColor(0xF0, 0xF3, 0xF7),
    "white":    RGBColor(0xFF, 0xFF, 0xFF),
    "cover_bg": RGBColor(0x1E, 0x2A, 0x38),
    "cover_sub":RGBColor(0x8A, 0x9B, 0xAE),
    "border":   RGBColor(0xDD, 0xDD, 0xDD),
}

# ══════════════════════════════════════════
#  HELPERS
# ══════════════════════════════════════════

prs = None
_page = [0]

def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = CLR["bg"]
    _page[0] += 1
    return s

def tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

def para(tf, text, size=Pt(18), color=None, bold=False, align=None, spacing=Pt(28), first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.font.size = size
    p.font.color.rgb = color or CLR["body"]
    p.font.bold = bold
    if align: p.alignment = align
    p.line_spacing = spacing
    return p

def bullet(tf, text, size=Pt(17), color=None):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = size
    p.font.color.rgb = color or CLR["body"]
    p.level = 0
    p.line_spacing = Pt(32)
    return p

def page_number(slide):
    t = tb(slide, 12.3, 7.05, 0.8, 0.35)
    para(t.text_frame, str(_page[0]), Pt(10), CLR["muted"], align=PP_ALIGN.RIGHT, first=True)

def title_bar(slide, text, subtitle=None):
    line = slide.shapes.add_shape(1, Inches(0.8), Inches(0.55), Inches(0.8), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = CLR["accent"]
    line.line.fill.background()
    t = tb(slide, 0.8, 0.15, 11.5, 0.6)
    para(t.text_frame, text, Pt(30), CLR["dark"], bold=True, first=True)
    if subtitle:
        t2 = tb(slide, 0.8, 0.7, 11.5, 0.4)
        para(t2.text_frame, subtitle, Pt(14), CLR["muted"], first=True)

def body_box(slide, top=1.2, left=0.8, width=11.5, height=5.5):
    t = tb(slide, left, top, width, height)
    t.text_frame.word_wrap = True
    return t.text_frame

def add_image(slide, img_path, left, top, width, height):
    if not os.path.exists(img_path):
        print(f"  [!] Missing: {img_path}")
        return None
    pic = slide.shapes.add_picture(img_path, Inches(left), Inches(top))
    max_w, max_h = Inches(width), Inches(height)
    ratio = min(max_w / pic.width, max_h / pic.height)
    if ratio < 1:
        pic.width  = int(pic.width  * ratio)
        pic.height = int(pic.height * ratio)
    pic.left = int(Inches(left) + (max_w - pic.width)  / 2)
    pic.top  = int(Inches(top)  + (max_h - pic.height) / 2)
    pic.line.color.rgb = CLR["border"]
    pic.line.width = Pt(0.5)
    return pic

def kpi_card(slide, left, top, value, label, color=None):
    w, h = 2.5, 1.3
    card = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = CLR["light"]
    card.line.fill.background()
    t = tb(slide, left + 0.2, top + 0.2, w - 0.4, 0.7)
    para(t.text_frame, value, Pt(36), color or CLR["accent"], bold=True, first=True)
    t2 = tb(slide, left + 0.2, top + 0.85, w - 0.4, 0.35)
    para(t2.text_frame, label, Pt(11), CLR["muted"], first=True)

# ══════════════════════════════════════════
#  MARKDOWN PARSER
# ══════════════════════════════════════════

IMG_RE = re.compile(r'!\[.*?\]\(\./(.+\.(?:png|jpg|jpeg|gif|webp|svg))\)', re.IGNORECASE)
KPI_RE = re.compile(r'(\d[\d,]*[KMB]?\+?\s*(?:stars?|forks?|万|亿|[KMB]?\s*USD|平台|[^a-zA-Z\s]{0,3}))')

def detect_kpis(text):
    """Detect key metrics in text. Returns list of (value, label) tuples."""
    kpis = []
    patterns = [
        (r'(\d[\d,\.]*[Kk]?\+?)\s*(?:的|个)?\s*(?:star|stars?)', 'GitHub Stars'),
        (r'(\d[\d,]*)\s*(?:的|个)?\s*(?:fork|forks?)', 'Forks'),
        (r'累计[^0-9]*?约\s*(\d[\d.]*[Kk]\s*USD)', '累计筹款'),
        (r'(\d[\d,]*\+?)\s*(?:个?\s*平台|平台覆盖)', '平台覆盖'),
        (r'(\d[\d,]*[Kk]?\s*(?:USD|usd))\s*(?:目标|已达成)', '筹款目标'),
        (r'(\d[\d,]*[Kk]\s*rep)', 'StackOverflow Rep'),
    ]
    seen = set()
    for pat, label in patterns:
        for m in re.finditer(pat, text, re.IGNORECASE):
            val = m.group(1).strip()
            if val not in seen:
                seen.add(val)
                kpis.append((val, label))
    return kpis[:4]

def clean_bullet(text):
    """Clean a single bullet text: remove URLs, zero-width chars, markdown artifacts."""
    # Remove zero-width characters
    text = re.sub(r'[​‌‍‎‏﻿]', '', text)
    # Remove bare URLs
    text = re.sub(r'https?://\S+', '', text)
    # Remove markdown links, keep text: [text](url) -> text
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
    # Remove markdown bold/italic markers
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
    text = re.sub(r'\*([^*]+)\*', r'\1', text)
    # Collapse whitespace
    text = re.sub(r'\s+', ' ', text).strip()
    return text

def find_assets_dir(md_path):
    """Auto-detect assets directory from markdown path."""
    base = md_path.parent
    md_name = md_path.stem
    candidates = [
        base / f"{md_name}.assets",
        base / "assets",
        base / f"{md_name}_assets",
    ]
    for c in candidates:
        if c.is_dir():
            return c
    return None

def resolve_image(img_rel, md_path, assets_dir):
    """Resolve an image reference to a full path."""
    # Try relative to md file
    full = md_path.parent / img_rel
    if full.exists():
        return str(full)
    # Try relative to assets dir
    if assets_dir:
        full = assets_dir / Path(img_rel).name
        if full.exists():
            return str(full)
    return None

def parse_markdown(md_text):
    """Parse markdown into structured sections.

    Returns: {
        'title': str,
        'sections': [{'title': str, 'subsections': [{'title': str, 'body': str, 'images': [str], 'kpis': [(v,l)]}]}]
    }
    """
    # Extract title from H1
    h1_m = re.match(r'^# (.+)', md_text)
    title = h1_m.group(1).strip() if h1_m else "Report"

    # Remove H1 from body
    body = re.sub(r'^# .*\n+', '', md_text)

    result = {'title': title, 'sections': []}

    # Split by ## headings
    sections = re.split(r'\n(?=## )', body.strip())
    for sec in sections:
        sec = sec.strip()
        if not sec:
            continue
        m = re.match(r'^## (.+)', sec)
        if not m:
            continue
        sec_title = m.group(1).strip()
        sec_body = sec[m.end():].strip()

        sec_data = {'title': sec_title, 'subsections': [], 'images': [], 'kpis': []}

        # Extract images from entire section
        sec_data['images'] = IMG_RE.findall(sec_body)
        sec_data['kpis'] = detect_kpis(sec_body)

        # Handle subsections (###)
        if '###' in sec_body:
            subs = re.split(r'\n(?=### )', sec_body)
            for sub in subs:
                sub = sub.strip()
                if not sub:
                    continue
                sm = re.match(r'^### (.+)', sub)
                if sm:
                    sub_title = sm.group(1).strip()
                    sub_body = sub[sm.end():].strip()
                else:
                    sub_title = ""
                    sub_body = sub

                sub_imgs = IMG_RE.findall(sub_body)
                sub_body_clean = IMG_RE.sub('', sub_body).strip()
                sub_body_clean = re.sub(r'\n{3,}', '\n\n', sub_body_clean)
                sub_kpis = detect_kpis(sub_body)

                sec_data['subsections'].append({
                    'title': sub_title,
                    'body': sub_body_clean,
                    'images': sub_imgs,
                    'kpis': sub_kpis,
                })
        else:
            # No subsections: treat entire body as one
            body_clean = IMG_RE.sub('', sec_body).strip()
            body_clean = re.sub(r'\n{3,}', '\n\n', body_clean)
            sec_data['subsections'].append({
                'title': '',
                'body': body_clean,
                'images': [],
                'kpis': [],
            })

        result['sections'].append(sec_data)

    return result

def extract_bullets(text, max_per_slide=6):
    """Extract bullet points from text. Returns list of cleaned strings."""
    bullets = []
    for line in text.split('\n'):
        line = line.strip()
        if not line:
            continue
        # Remove markdown list markers
        line = re.sub(r'^[-*]\s+', '', line)
        # Apply full cleaning
        line = clean_bullet(line)
        if not line:
            continue
        # Trim long lines
        if len(line) > 120:
            # Try to split at sentence boundaries
            parts = re.split(r'(?<=[。；])', line)
            for p in parts:
                p = p.strip()
                if p:
                    bullets.append(p)
        else:
            bullets.append(line)
    return bullets[:max_per_slide]

def chunk_bullets(bullets, max_per_slide=6):
    """Split bullets into slide-sized chunks."""
    chunks = []
    for i in range(0, len(bullets), max_per_slide):
        chunks.append(bullets[i:i+max_per_slide])
    return chunks

# ══════════════════════════════════════════
#  AUTO SLIDE GENERATION
# ══════════════════════════════════════════

def gen_cover(title, date_str=None):
    """Generate cover slide."""
    s = new_slide()
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = CLR["cover_bg"]

    bar = s.shapes.add_shape(1, Inches(1.5), Inches(3.0), Inches(0.08), Inches(1.6))
    bar.fill.solid(); bar.fill.fore_color.rgb = CLR["accent"]
    bar.line.fill.background()

    # Split title into main + subtitle on first space or delimiter
    parts = title.split(' — ', 1) if ' — ' in title else title.split('：', 1)
    main_title = parts[0].strip()
    sub_title = parts[1].strip() if len(parts) > 1 else ""

    t = tb(s, 1.5, 2.8, 10.0, 1.4)
    tf = t.text_frame; tf.word_wrap = True
    para(tf, main_title, Pt(48), CLR["white"], bold=True, first=True)
    if sub_title:
        para(tf, sub_title, Pt(32), CLR["accent"], bold=True, spacing=Pt(48))
    else:
        para(tf, "", Pt(8), spacing=Pt(8))

    t2 = tb(s, 1.5, 5.0, 10.0, 0.6)
    para(t2.text_frame, "舆情调查报告" if "舆情" in main_title else "调查报告", Pt(22), CLR["cover_sub"], first=True)
    t3 = tb(s, 1.5, 5.6, 10.0, 0.4)
    para(t3.text_frame, date_str or "", Pt(14), CLR["cover_sub"], first=True)


def gen_section_divider(section_title):
    """Generate a section divider slide."""
    s = new_slide()
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = CLR["cover_bg"]

    bar = s.shapes.add_shape(1, Inches(1.5), Inches(3.2), Inches(0.08), Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = CLR["accent"]
    bar.line.fill.background()

    t = tb(s, 1.5, 3.0, 10.0, 1.0)
    para(t.text_frame, section_title, Pt(36), CLR["white"], bold=True, first=True)


def gen_content_text(slide_title, bullets_list, subtitle=None):
    """Generate a text-only content slide."""
    s = new_slide()
    title_bar(s, slide_title, subtitle)
    tf = body_box(s, width=11.5, height=5.8)
    para(tf, "", Pt(6), first=True)
    for b in bullets_list:
        bullet(tf, f"  {b}")
    page_number(s)
    return s


def gen_content_with_image(slide_title, bullets_list, img_path, subtitle=None):
    """Generate a content slide with text left + image right."""
    s = new_slide()
    title_bar(s, slide_title, subtitle)
    tf = body_box(s, width=5.5)
    para(tf, "", Pt(6), first=True)
    for b in bullets_list:
        bullet(tf, f"  {b}")
    add_image(s, img_path, 7.0, 1.3, 5.6, 5.0)
    page_number(s)
    return s


def gen_summary(kpis, bullets_list):
    """Generate summary slide with KPI cards."""
    s = new_slide()
    title_bar(s, "摘要", "Executive Summary")

    # KPI cards in a row
    if kpis:
        positions = [0.8, 3.6, 6.4, 9.2]
        for i, (value, label) in enumerate(kpis[:4]):
            kpi_card(s, positions[i], 1.3, value, label)

    # Bullet points below KPIs
    tf = body_box(s, top=2.9, height=4.0)
    para(tf, "", Pt(6), first=True)
    for b in bullets_list[:4]:
        bullet(tf, f"  {b}")

    page_number(s)
    return s


def gen_conclusion(title, bullets_list):
    """Generate conclusion slide with dark background."""
    s = new_slide()
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = CLR["cover_bg"]

    bar = s.shapes.add_shape(1, Inches(1.5), Inches(2.2), Inches(0.08), Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = CLR["accent"]
    bar.line.fill.background()

    t = tb(s, 1.5, 2.0, 10.0, 0.8)
    para(t.text_frame, title, Pt(36), CLR["white"], bold=True, first=True)

    t2 = tb(s, 1.5, 3.5, 10.0, 3.0)
    tf2 = t2.text_frame; tf2.word_wrap = True
    first = True
    for pt in bullets_list[:6]:
        if pt:
            para(tf2, f"  {pt}", Pt(17), CLR["white"], first=first, spacing=Pt(30))
        first = False

    page_number(s)


def auto_generate(md_path, output_path, title_override=None, date_str=None, assets_dir=None):
    """Parse markdown and auto-generate PPTX slides.

    Args:
        md_path: Path to the markdown report
        output_path: Path for the output PPTX
        title_override: Optional custom title
        assets_dir: Optional custom assets directory path
    """
    global prs
    md_path = Path(md_path)
    md_text = md_path.read_text(encoding="utf-8")

    if assets_dir is None:
        assets_dir = find_assets_dir(md_path)

    def img_resolve(name):
        """Resolve an image filename to full path."""
        if assets_dir:
            full = assets_dir / Path(name).name
            if full.exists():
                return str(full)
        full = md_path.parent / name
        if full.exists():
            return str(full)
        return str(md_path.parent / name)

    parsed = parse_markdown(md_text)
    report_title = title_override or parsed['title']

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slides = []

    # 1. Cover
    slides.append(('cover', report_title, date_str))

    # 2. Gather all KPIs and summary bullets
    all_kpis = []
    all_summary_bullets = []
    for sec in parsed['sections']:
        all_kpis.extend(sec.get('kpis', []))
        for sub in sec.get('subsections', []):
            all_kpis.extend(sub.get('kpis', []))
            bullets_text = extract_bullets(sub['body'])
            all_summary_bullets.extend(bullets_text[:2])

    # Deduplicate KPIs
    seen_kpis = set()
    unique_kpis = []
    for v, l in all_kpis:
        if v not in seen_kpis:
            seen_kpis.add(v)
            unique_kpis.append((v, l))
    unique_kpis = unique_kpis[:4]

    # Summary slide
    if unique_kpis or all_summary_bullets:
        # Use first section's first subsection bullets as summary
        first_section = parsed['sections'][0] if parsed['sections'] else None
        summary_bullets = []
        if first_section:
            for sub in first_section['subsections'][:1]:
                summary_bullets = extract_bullets(sub['body'])[:4]
        if not summary_bullets:
            summary_bullets = all_summary_bullets[:4]

        slides.append(('summary', unique_kpis, summary_bullets))

    # 3. Sections — merge subsections into section-level slides
    for i, sec in enumerate(parsed['sections']):
        sec_title = sec['title']

        # Collect bullets and images from all subsections in this section
        sec_bullets = []
        sec_images = []
        first_subtitle = None
        for sub in sec['subsections']:
            if sub['title'] and first_subtitle is None:
                first_subtitle = sub['title']
            sec_bullets.extend(extract_bullets(sub['body']))
            for img in sub['images']:
                path = img_resolve(img)
                if os.path.exists(path) and path not in sec_images:
                    sec_images.append(path)

        if not sec_bullets and not sec_images:
            continue

        # Chunk bullets: 5 per slide, at most 1 image per slide
        MAX_PER = 5
        bullet_chunks = chunk_bullets(sec_bullets, MAX_PER) if sec_bullets else [[]]
        img_slice = 0

        for ci, chunk in enumerate(bullet_chunks):
            if not chunk and img_slice >= len(sec_images):
                continue
            stitle = sec_title if ci == 0 else f"{sec_title} (续{ci + 1})"
            img = sec_images[img_slice] if ci < len(sec_images) else None
            if img:
                img_slice += 1
            sub = first_subtitle if ci == 0 else None
            if img and chunk:
                slides.append(('content_image', stitle, chunk, img, sub))
            elif chunk:
                slides.append(('content_text', stitle, chunk, sub))

    # Generate all slides
    for slide_spec in slides:
        spec_type = slide_spec[0]

        if spec_type == 'cover':
            gen_cover(slide_spec[1], slide_spec[2])
        elif spec_type == 'summary':
            gen_summary(slide_spec[1], slide_spec[2])
        elif spec_type == 'content_text':
            gen_content_text(slide_spec[1], slide_spec[2], slide_spec[3] if len(slide_spec) > 3 else None)
        elif spec_type == 'content_image':
            gen_content_with_image(slide_spec[1], slide_spec[2], slide_spec[3], slide_spec[4] if len(slide_spec) > 4 else None)

    # 5. Conclusion slide from last section
    if parsed['sections']:
        last_sec = parsed['sections'][-1]
        conclusion_bullets = []
        for sub in last_sec['subsections']:
            conclusion_bullets.extend(extract_bullets(sub['body']))
        if conclusion_bullets:
            gen_conclusion(last_sec['title'], conclusion_bullets[:6])

    # Save
    prs.save(str(output_path))
    return len(prs.slides), output_path


# ══════════════════════════════════════════
#  OUTLINE MODE — parse .slide.md (Marp format) and generate PPTX
# ══════════════════════════════════════════

def parse_outline(text):
    """Parse Marp-format .slide.md into a list of slide specs.

    Each spec dict:
      type: 'cover' | 'lead' | 'content' | 'content_image'
      title: str
      subtitle: str
      bullets: [(text, bold_parts)]
      paragraphs: [str]
      images: [(path, placement)]  # placement: 'right' | 'full' | 'inline'
      kpis: [(value, label)]
      image_right: str or None  # path for right-side image
    """
    # Remove YAML frontmatter
    text = re.sub(r'^---\n.*?\n---\n', '', text, flags=re.DOTALL)

    # Split by slide separator
    raw_slides = re.split(r'\n---\n', text)
    slides = []
    is_first = True

    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue

        spec = {
            'type': 'content',
            'title': '',
            'subtitle': '',
            'bullets': [],
            'paragraphs': [],
            'images': [],
            'kpis': [],
            'image_right': None,
        }

        # Detect lead class
        is_lead = '<!-- _class: lead -->' in raw
        raw = raw.replace('<!-- _class: lead -->', '').strip()

        # Extract image directives: ![bg right:40%](path), ![bg](path), ![alt](path)
        img_re = re.compile(r'!\[(.*?)\]\((.+?)\)')
        for m in img_re.finditer(raw):
            alt = m.group(1).strip()
            path = m.group(2).strip()
            # Remove ./ prefix
            if path.startswith('./'):
                path = path[2:]
            # URL-decode the path (Marp encodes parens and spaces)
            from urllib.parse import unquote
            path = unquote(path)
            if alt.startswith('bg right'):
                spec['image_right'] = path
                spec['type'] = 'content_image'
            elif alt.startswith('bg'):
                spec['type'] = 'content_image'
                spec['image_right'] = path
            else:
                spec['images'].append((path, 'inline'))
        # Remove image lines
        raw = img_re.sub('', raw)

        # Parse remaining content line by line
        lines = raw.strip().split('\n')
        for line in lines:
            line = line.strip()
            if not line:
                continue

            # Headings
            if line.startswith('# ') and not spec['title']:
                spec['title'] = line[2:].strip()
                # Strip markdown formatting from title
                spec['title'] = re.sub(r'\*\*([^*]+)\*\*', r'\1', spec['title'])
                spec['title'] = re.sub(r'`([^`]+)`', r'\1', spec['title'])
            elif line.startswith('## '):
                s = line[3:].strip()
                s = re.sub(r'\*\*([^*]+)\*\*', r'\1', s)
                s = re.sub(r'`([^`]+)`', r'\1', s)
                if not spec['subtitle']:
                    spec['subtitle'] = s
                else:
                    spec['paragraphs'].append(s)
            elif line.startswith('- '):
                bullet_text = line[2:].strip()
                spec['bullets'].append(bullet_text)
            else:
                spec['paragraphs'].append(line)

        # Detect KPIs from bold text in paragraphs and bullets
        all_text = ' '.join(spec['paragraphs'] + spec['bullets'] + [spec['title'], spec['subtitle']])
        kpi_patterns = [
            (r'\*\*([\d,]+[Kk]?\+?\s*stars?)\s*/\s*([\d,]+\s*forks?)\*\*', 'Stars / Forks'),
            (r'\*\*([\d,.]+[KkMm]?\s*USD)\*\*', 'Funding'),
            (r'\*\*([\d,]+[Kk]?\s*rep)\*\*', 'Reputation'),
            (r'\*\*([\d,]+[Kk]?\+?\s*stars?)\*\*', 'GitHub Stars'),
            (r'\*\*([\d,]+\s*forks?)\*\*', 'Forks'),
            (r'\*\*([\d,]+\+?\s*(?:平台|Platforms?|个?\s*平台))\*\*', 'Platforms'),
        ]
        for pat, label in kpi_patterns:
            for km in re.finditer(pat, all_text, re.I):
                groups = km.groups()
                if len(groups) == 2 and '/' in km.group(0):
                    # Multi-value KPI like "3,000+ stars / 290 forks"
                    spec['kpis'].append((groups[0].strip(), 'Stars'))
                    spec['kpis'].append((groups[1].strip(), 'Forks'))
                else:
                    spec['kpis'].append((groups[0].strip(), label))
        # Deduplicate
        seen = set()
        unique_kpis = []
        for v, l in spec['kpis']:
            if v not in seen:
                seen.add(v)
                unique_kpis.append((v, l))
        spec['kpis'] = unique_kpis[:4]

        # Determine slide type
        if is_lead:
            spec['type'] = 'cover' if is_first else 'lead'
        elif spec['image_right'] and spec['bullets']:
            spec['type'] = 'content_image'
        elif spec['image_right']:
            spec['type'] = 'content_image'

        is_first = False
        slides.append(spec)

    return slides


def gen_from_outline(slides, output_path, assets_dir=None):
    """Generate PPTX from parsed outline specs using v2 design system."""
    global prs
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    if assets_dir is None:
        assets_dir = Path(output_path).parent

    def resolve_img(name):
        """Resolve image filename to full path."""
        if not name:
            return None
        # Try assets dir first
        if assets_dir:
            # Check subdirectories for .assets folder
            for d in [assets_dir] + list(assets_dir.glob("*.assets")):
                full = Path(d) / Path(name).name
                if full.exists():
                    return str(full)
        # Try the name as-is and relative to output
        for base in [Path('.'), Path(output_path).parent]:
            full = base / name
            if full.exists():
                return str(full)
        return None

    _page[0] = 0

    for spec in slides:
        stype = spec['type']

        if stype == 'cover':
            s = new_slide()
            s.background.fill.solid()
            s.background.fill.fore_color.rgb = CLR["cover_bg"]

            bar = s.shapes.add_shape(1, Inches(1.5), Inches(3.0), Inches(0.08), Inches(1.6))
            bar.fill.solid(); bar.fill.fore_color.rgb = CLR["accent"]
            bar.line.fill.background()

            t = tb(s, 1.5, 2.8, 10.0, 1.4)
            tf = t.text_frame; tf.word_wrap = True
            title = spec['title'] or "Report"
            subtitle = spec['subtitle'] or ""
            para(tf, title, Pt(48), CLR["white"], bold=True, first=True)
            if subtitle:
                para(tf, subtitle, Pt(32), CLR["accent"], bold=True, spacing=Pt(48))

            # Report type line from first paragraph
            report_type = spec['paragraphs'][0] if spec['paragraphs'] else "调查报告"
            t2 = tb(s, 1.5, 5.0, 10.0, 0.6)
            para(t2.text_frame, report_type, Pt(22), CLR["cover_sub"], first=True)
            t3 = tb(s, 1.5, 5.6, 10.0, 0.4)
            para(t3.text_frame, "", Pt(14), CLR["cover_sub"], first=True)

        elif stype == 'lead':
            # Section divider — dark background
            s = new_slide()
            s.background.fill.solid()
            s.background.fill.fore_color.rgb = CLR["cover_bg"]

            bar = s.shapes.add_shape(1, Inches(1.5), Inches(3.2), Inches(0.08), Inches(1.0))
            bar.fill.solid(); bar.fill.fore_color.rgb = CLR["accent"]
            bar.line.fill.background()

            t = tb(s, 1.5, 3.0, 10.0, 1.5)
            tf = t.text_frame; tf.word_wrap = True
            para(tf, spec['title'] or spec['subtitle'], Pt(36), CLR["white"], bold=True, first=True)
            for p_text in spec['paragraphs'][:3]:
                clean = re.sub(r'\*\*([^*]+)\*\*', r'\1', p_text)
                clean = re.sub(r'`([^`]+)`', r'\1', clean)
                para(tf, clean, Pt(18), CLR["cover_sub"])
            for b_text in spec['bullets'][:3]:
                clean = re.sub(r'\*\*([^*]+)\*\*', r'\1', b_text)
                clean = re.sub(r'`([^`]+)`', r'\1', clean)
                bullet(tf, f"  {clean}")

        elif stype in ('content', 'content_image'):
            s = new_slide()
            title_bar(s, spec['title'] or spec['subtitle'], spec['subtitle'] if spec['title'] else None)

            # KPI cards at top if present
            if spec['kpis']:
                positions = [0.8, 3.6, 6.4, 9.2]
                for i, (value, label) in enumerate(spec['kpis'][:4]):
                    kpi_card(s, positions[i], 1.3, value, label)
                body_top = 2.9
                body_height = 4.0
            else:
                body_top = 1.2
                body_height = 5.5

            # Determine layout: text width depends on image
            has_image = bool(spec['image_right'])
            text_width = 5.5 if has_image else 11.5

            tf = body_box(s, top=body_top, width=text_width, height=body_height)
            first = True

            # Render paragraphs
            for p_text in spec['paragraphs']:
                if first:
                    para(tf, "", Pt(6), first=True)
                    first = False
                # Clean markdown formatting
                clean = re.sub(r'\*\*([^*]+)\*\*', r'\1', p_text)
                clean = re.sub(r'`([^`]+)`', r'\1', clean)
                # Handle ### sub-headers
                if clean.startswith('### '):
                    clean = clean[4:].strip()
                    para(tf, "", Pt(4), spacing=Pt(4))
                    para(tf, clean, Pt(20), CLR["accent"], bold=True, spacing=Pt(24))
                # Handle table rows
                elif clean.startswith('|') and clean.endswith('|'):
                    clean = clean.strip('|')
                    cells = [c.strip() for c in clean.split('|')]
                    # Skip separator rows like "---|---"
                    if all(re.match(r'^[-:]+$', c) for c in cells):
                        continue
                    # Render as indented key-value
                    if len(cells) >= 2:
                        para(tf, f"  {cells[0]}：{cells[1]}", Pt(14), CLR["body"], spacing=Pt(24))
                else:
                    para(tf, clean, Pt(17), CLR["body"], spacing=Pt(28))

            # Render bullets
            for b_text in spec['bullets']:
                clean = re.sub(r'\*\*([^*]+)\*\*', r'\1', b_text)
                clean = re.sub(r'`([^`]+)`', r'\1', clean)
                # Handle ### sub-headers in bullet text
                if clean.startswith('### '):
                    if not first:
                        para(tf, "", Pt(4), spacing=Pt(4))
                    clean = clean[4:].strip()
                    para(tf, clean, Pt(20), CLR["accent"], bold=True, spacing=Pt(24))
                    first = False
                else:
                    if first:
                        para(tf, "", Pt(6), first=True)
                        first = False
                    bullet(tf, f"  {clean}")

            # Image on right side
            if has_image:
                img_path = resolve_img(spec['image_right'])
                if img_path and os.path.exists(img_path):
                    add_image(s, img_path, 7.0, body_top, 5.6, body_height)

            page_number(s)

    prs.save(str(output_path))
    return len(prs.slides), output_path


# ══════════════════════════════════════════
#  MANUAL MODE — edit slide functions below
# ══════════════════════════════════════════

def slide01_cover():
    """Cover slide."""
    s = new_slide()
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = CLR["cover_bg"]

    bar = s.shapes.add_shape(1, Inches(1.5), Inches(3.0), Inches(0.08), Inches(1.6))
    bar.fill.solid(); bar.fill.fore_color.rgb = CLR["accent"]
    bar.line.fill.background()

    t = tb(s, 1.5, 2.8, 10.0, 1.4)
    tf = t.text_frame; tf.word_wrap = True
    para(tf, "Main Title", Pt(54), CLR["white"], bold=True, first=True)
    para(tf, "Subtitle", Pt(36), CLR["accent"], bold=True, spacing=Pt(48))

    t2 = tb(s, 1.5, 5.0, 10.0, 0.6)
    para(t2.text_frame, "Report Type", Pt(22), CLR["cover_sub"], first=True)
    t3 = tb(s, 1.5, 5.6, 10.0, 0.4)
    para(t3.text_frame, "2026 年 6 月", Pt(14), CLR["cover_sub"], first=True)


def slide02_example():
    """Example content slide."""
    s = new_slide()
    title_bar(s, "Slide Title", "Optional Subtitle")

    tf = body_box(s, width=5.5)
    para(tf, "", Pt(6), first=True)
    bullet(tf, "  First bullet point")
    bullet(tf, "  Second bullet point")
    bullet(tf, "  Third bullet point")

    add_image(s, "example.png", 7.0, 1.3, 5.6, 5.0)
    page_number(s)


# ══════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════

def main():
    global prs

    # Parse command line: find .md input and optional -o/-t flags
    argv = sys.argv[1:]
    md_path = None
    output_path = None
    title_override = None

    i = 0
    while i < len(argv):
        arg = argv[i]
        if arg in ('-o', '--output') and i + 1 < len(argv):
            i += 1
            output_path = Path(argv[i])
        elif arg.startswith('--output='):
            output_path = Path(arg.split('=', 1)[1])
        elif arg.startswith('-o='):
            output_path = Path(arg[3:])
        elif arg in ('-t', '--title') and i + 1 < len(argv):
            i += 1
            title_override = argv[i]
        elif arg.startswith('--title='):
            title_override = arg.split('=', 1)[1]
        elif arg.startswith('-t='):
            title_override = arg[3:]
        elif arg.endswith('.md'):
            md_path = Path(arg)
        i += 1

    if md_path and md_path.exists():
        if output_path is None:
            output_path = md_path.parent / f"{md_path.stem}.pptx"

        # Detect mode: .slide.md → outline, .md → auto
        if md_path.suffix == '.md' and (md_path.stem.endswith('.slide') or 'slide' in md_path.name.lower()):
            # ── Outline mode (.slide.md) ──
            print(f"Mode:   Outline (Marp .slide.md)")
            print(f"Input:  {md_path}")
            print(f"Output: {output_path}")
            print()

            text = md_path.read_text(encoding='utf-8')
            specs = parse_outline(text)
            print(f"Parsed {len(specs)} slide specs from outline")

            assets_dir = find_assets_dir(md_path)
            n_slides, out = gen_from_outline(specs, output_path, assets_dir)
            size_mb = output_path.stat().st_size / 1024 / 1024
            print(f"OK: {n_slides} slides, {size_mb:.1f} MB -> {output_path}")

        else:
            # ── Auto mode (.md) ──
            import datetime
            mtime = md_path.stat().st_mtime
            dt = datetime.datetime.fromtimestamp(mtime)
            date_str = f"{dt.year} 年 {dt.month} 月"

            print(f"Mode:   Auto (raw .md)")
            print(f"Input:  {md_path}")
            print(f"Output: {output_path}")
            if title_override:
                print(f"Title:  {title_override}")
            print()

            n_slides, out = auto_generate(md_path, output_path, title_override, date_str)
            size_mb = output_path.stat().st_size / 1024 / 1024
            print(f"OK: {n_slides} slides, {size_mb:.1f} MB -> {output_path}")

    else:
        # ── Manual mode ──
        prs = Presentation()
        prs.slide_width  = W
        prs.slide_height = H

        slides = [
            slide01_cover,
            slide02_example,
            # Add more slide functions here...
        ]

        for fn in slides:
            fn()

        output = Path(__file__).parent / "output.pptx"
        prs.save(str(output))
        size_mb = output.stat().st_size / 1024 / 1024
        print(f"OK: {len(slides)} slides, {size_mb:.1f} MB -> {output}")

if __name__ == "__main__":
    main()
